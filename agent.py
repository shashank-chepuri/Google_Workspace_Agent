#!/usr/bin/env python3
import os
import re
import base64
import json
import io
import string
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload
from google.oauth2.credentials import Credentials
from email.mime.multipart import MIMEMultipart
from email.mime.base import MIMEBase
from email.mime.text import MIMEText
from email import encoders
import cohere
from dotenv import load_dotenv

# Optional parsers
try:
    import PyPDF2
except ImportError:
    PyPDF2 = None
try:
    import docx
except ImportError:
    docx = None
try:
    from pptx import Presentation
except ImportError:
    Presentation = None

# Load environment variables
load_dotenv()
cohere_api_key = os.getenv("COHERE_API_KEY")
if not cohere_api_key:
    raise ValueError("‚ùå COHERE_API_KEY not found in .env file")

co = cohere.Client(cohere_api_key)

# --------------------
# CONFIG
# --------------------
SCOPES = [
    "https://www.googleapis.com/auth/drive.readonly",
    "https://www.googleapis.com/auth/spreadsheets.readonly",
    "https://www.googleapis.com/auth/gmail.send",
    "https://www.googleapis.com/auth/documents.readonly",
    "https://www.googleapis.com/auth/presentations.readonly"
]

# --------------------
# Draft state
# --------------------
current_draft = {
    "subject": None,
    "body": None,
    "recipients": [],
    "context": ""  # Store original context for reference
}

# --------------------
# Utility Functions
# --------------------
def _escape_for_drive_q(s: str) -> str:
    return s.replace("'", "\\'").strip()

def _print_files(files):
    if not files:
        print("No files found.")
        return
    print("\nüìÇ Matching files:")
    print("-" * 100)
    for f in files:
        mid = f.get("id", "")
        name = f.get("name", "")
        mtype = f.get("mimeType", "")
        print(f"{name}  |  ID: {mid}  |  {mtype}")
    print("-" * 100)

def _unique_by_id(files):
    seen = set()
    out = []
    for f in files:
        fid = f.get("id")
        if fid and fid not in seen:
            seen.add(fid)
            out.append(f)
    return out

def _clean_text(s: str) -> str:
    if not s:
        return s
    allowed = set(string.printable) | set("‚Äú‚Äù‚Äò‚Äô‚Äî‚Äì‚Ä¢¬∑\u200b\t\n\r")
    cleaned = ''.join(ch if (ch in allowed or ch.isprintable()) else ' ' for ch in s)
    cleaned = cleaned.replace('\u200b', '')
    cleaned = re.sub(r'[ \t]+', ' ', cleaned)
    cleaned = re.sub(r'\n{3,}', '\n\n', cleaned)
    return cleaned.strip()

def _likely_binary(b: bytes) -> bool:
    if not b:
        return False
    if b.count(b'\x00') > 0:
        return True
    sample = b[:2048]
    try:
        s = sample.decode('utf-8', errors='ignore')
    except Exception:
        return True
    printable_ratio = sum(ch.isprintable() or ch in '\n\r\t' for ch in s) / max(1, len(s))
    return printable_ratio < 0.7

def parse_recipients(text: str):
    """Extract email addresses from text."""
    # Improved regex for better email matching
    emails = re.findall(r'[\w\.-]+@[\w\.-]+\.\w+', text)
    out = []
    seen = set()
    for e in emails:
        e = e.strip().strip(",;")
        if e and e not in seen:
            seen.add(e)
            out.append(e)
    return out

def extract_context_after_keyword(command: str, keywords: list):
    """Extract everything after the first occurrence of any keyword."""
    command_lower = command.lower()
    for kw in keywords:
        if kw in command_lower:
            idx = command_lower.find(kw) + len(kw)
            return command[idx:].strip()
    return ""

# --------------------
# AI Command Parser (Improved)
# --------------------
def parse_command(command: str):
    """
    Improved command parser with better draft email handling.
    """
    # Try LLM primary path first
    prompt = f"""
You are a command interpreter for a Google Drive + Email assistant.
Return ONLY a single valid JSON object, no prose, no markdown.

Valid actions:
1. {{"action":"list_files"}}
2. {{"action":"search_files","keyword":"budget"}}
3. {{"action":"send_file","file_name":["file1.pdf"],"email":"user@example.com"}}
4. {{"action":"summarize_file","file_name":"report.pdf"}} OR {{"action":"summarize_file","file_id":"123"}}
5. {{"action":"draft_email","text":"I need to request sick leave for 2 days starting Monday"}}
6. {{"action":"refine_draft","instruction":"make it more formal"}}
7. {{"action":"send_draft","email":["recipient@example.com"]}}
8. {{"action":"show_draft"}}
9. {{"action":"clear_draft"}}
10. {{"action":"help"}}
11. {{"action":"exit"}}

Parse this command and output JSON only:
Command: {command}
"""

    try:
        resp = co.chat(model="command-r-plus", message=prompt, temperature=0)
        text = (resp.text or "").strip()
        
        # Clean markdown code blocks
        if text.startswith("```"):
            text = re.sub(r"^```(?:json)?\s*|\s*```$", "", text, flags=re.IGNORECASE | re.DOTALL).strip()
        
        data = json.loads(text)
        if isinstance(data, dict) and "action" in data:
            return data
    except Exception:
        pass  # Fall through to heuristics

    # ---------- IMPROVED HEURISTICS ----------
    low = command.lower().strip()
    orig = command.strip()
    
    # Exit/Help
    if re.search(r"\b(exit|quit|bye|close)\b", low):
        return {"action": "exit"}
    if re.search(r"\b(help|\?|what can you do|commands)\b", low):
        return {"action": "help"}
    
    # Draft management
    if re.search(r"\b(show|display|view)\b.*\bdraft\b", low):
        return {"action": "show_draft"}
    if re.search(r"\b(clear|delete|discard|erase)\b.*\bdraft\b", low):
        return {"action": "clear_draft"}
    
    # ===== IMPROVED DRAFT EMAIL PARSING =====
    # Pattern 1: "draft X" where X is the complete context
    draft_patterns = [
        r"^draft\s+(.+)$",  # starts with "draft"
        r"^compose\s+(.+)$",  # starts with "compose"
        r"^write\s+(.+)$",  # starts with "write"
        r"^create\s+(?:an?\s+)?email\s+(.+)$",  # "create email about X"
        r"^draft\s+(?:an?\s+)?(?:email|letter|message)\s+(.+)$",  # "draft an email about X"
    ]
    
    for pattern in draft_patterns:
        match = re.match(pattern, orig, re.IGNORECASE)
        if match:
            context = match.group(1).strip()
            # If context is just 1-2 words, treat as insufficient
            if len(context.split()) <= 2:
                return {"action": "draft_email", "text": ""}  # Will trigger interactive prompt
            return {"action": "draft_email", "text": context}
    
    # Pattern 2: "draft" with no context or minimal context
    if re.match(r"^draft\s*$", low) or low == "draft" or low.startswith("draft ") and len(low.split()) <= 2:
        return {"action": "draft_email", "text": ""}
    
    # Pattern 3: Extract context after "about" or "regarding"
    about_match = re.search(r"\b(draft|compose|write)\b.*?\b(?:about|regarding|for)\s+(.+)", orig, re.IGNORECASE)
    if about_match:
        return {"action": "draft_email", "text": about_match.group(2).strip()}
    
    # Refine draft
    refine_keywords = ["make it", "more formal", "more polite", "shorten", "shorter", 
                      "longer", "expand", "rewrite", "reword", "polite", "formal", 
                      "casual", "tone", "concise", "professional"]
    if any(kw in low for kw in refine_keywords):
        return {"action": "refine_draft", "instruction": orig}
    
    # Send draft
    send_draft_match = re.search(r"(send|email|mail)\s+(?:the\s+)?(draft|message|email)\s+to\s+(.+)", low)
    if send_draft_match:
        recipients_str = send_draft_match.group(3).strip()
        recipients = parse_recipients(recipients_str)
        return {"action": "send_draft", "email": recipients if recipients else [recipients_str]}
    
    # List files
    if re.search(r"\b(list|show).*files?\b", low):
        return {"action": "list_files"}
    
    # Search files
    search_match = re.search(r"\b(search|find|look\s*up)\b\s+(.*)", low)
    if search_match:
        rest = search_match.group(2).strip().strip("'\"")
        kws = [k.strip() for k in re.split(r"[,&]| and ", rest) if k.strip()]
        if kws:
            return {"action": "search_files", "keywords": kws if len(kws) > 1 else kws}
        return {"action": "search_files", "keyword": rest}
    
    # Summarize by ID
    id_match = re.search(r"(summari[sz]e|summary|tl;dr).*?id\s+([a-zA-Z0-9_\-]+)", low)
    if id_match:
        return {"action": "summarize_file", "file_id": id_match.group(2)}
    
    # Summarize by name
    summary_match = re.search(r"(summari[sz]e|summary|tl;dr|what'?s inside)\s+(?:file\s+)?(.+)", low)
    if summary_match:
        return {"action": "summarize_file", "file_name": summary_match.group(2).strip()}
    
    # Send files
    send_match = re.search(r"(send|email|mail)\s+(.+?)\s+to\s+([^\s]+(?:@[^\s]+)?.*?)$", low)
    if send_match:
        rest = send_match.group(2).strip()
        recipients = parse_recipients(send_match.group(3))
        tokens = [t.strip().strip("'\"") for t in re.split(r"[,&]| and ", rest) if t.strip()]
        
        # Check if this is about sending a draft
        if "draft" in rest or "message" in rest or "email" in rest:
            return {"action": "send_draft", "email": recipients if recipients else [send_match.group(3).strip()]}
        
        return {"action": "send_file", "file_name": tokens, "email": recipients[0] if recipients else send_match.group(3).strip()}
    
    return None

# --------------------
# Google API Auth
# --------------------
def get_services():
    try:
        creds = Credentials.from_authorized_user_file("token.json", SCOPES)
    except FileNotFoundError:
        raise FileNotFoundError("‚ùå token.json not found. Please run Google API auth flow.")
    
    drive_service = build("drive", "v3", credentials=creds)
    sheets_service = build("sheets", "v4", credentials=creds)
    docs_service = build("docs", "v1", credentials=creds)
    slides_service = build("slides", "v1", credentials=creds)
    gmail_service = build("gmail", "v1", credentials=creds)
    return drive_service, sheets_service, docs_service, slides_service, gmail_service

# --------------------
# Drive Operations
# --------------------
def list_all_files(drive_service):
    results = drive_service.files().list(
        pageSize=100,
        q="trashed = false",
        fields="files(id, name, mimeType, modifiedTime, shortcutDetails(targetId,targetMimeType))"
    ).execute()
    files = results.get("files", [])
    _print_files(files)

def search_files(drive_service, keyword):
    kw = _escape_for_drive_q(keyword)
    results = drive_service.files().list(
        q=f"trashed = false and name contains '{kw}'",
        fields="files(id, name, mimeType, modifiedTime, shortcutDetails(targetId,targetMimeType))"
    ).execute()
    return results.get("files", [])

def get_file_by_id(drive_service, file_id):
    return drive_service.files().get(
        fileId=file_id,
        fields="id,name,mimeType,modifiedTime,shortcutDetails(targetId,targetMimeType)"
    ).execute()

def _deref_shortcut(drive_service, file_dict):
    if not file_dict:
        return file_dict
    if file_dict.get("mimeType") == "application/vnd.google-apps.shortcut":
        target_id = (file_dict.get("shortcutDetails") or {}).get("targetId")
        if target_id:
            try:
                return get_file_by_id(drive_service, target_id)
            except Exception:
                pass
    return file_dict

# --------------------
# File Download & Content Extraction
# --------------------
def download_file(drive_service, file_dict):
    file_dict = _deref_shortcut(drive_service, file_dict)
    file_id = file_dict["id"]
    mime_type = file_dict["mimeType"]
    
    export_map = {
        "application/vnd.google-apps.document": "application/pdf",
        "application/vnd.google-apps.spreadsheet": "application/pdf",
        "application/vnd.google-apps.presentation": "application/pdf",
        "application/vnd.google-apps.drawing": "application/pdf",
    }
    
    if mime_type in export_map:
        request = drive_service.files().export_media(fileId=file_id, mimeType=export_map[mime_type])
    else:
        request = drive_service.files().get_media(fileId=file_id)
    
    fh = io.BytesIO()
    downloader = MediaIoBaseDownload(fh, request)
    done = False
    while not done:
        _, done = downloader.next_chunk()
    fh.seek(0)
    return fh.read()

def extract_file_content(drive_service, sheets_service, docs_service, slides_service, file_dict):
    file_dict = _deref_shortcut(drive_service, file_dict)
    file_id = file_dict["id"]
    mime_type = file_dict["mimeType"]
    
    try:
        # Google Docs
        if mime_type == "application/vnd.google-apps.document":
            doc = docs_service.documents().get(documentId=file_id).execute()
            text = []
            for c in doc.get("body", {}).get("content", []):
                if "paragraph" in c:
                    for e in c["paragraph"].get("elements", []):
                        tr = e.get("textRun", {})
                        if tr and "content" in tr:
                            text.append(tr["content"])
            return _clean_text("".join(text))
        
        # Google Sheets
        if mime_type == "application/vnd.google-apps.spreadsheet":
            sheet = sheets_service.spreadsheets().values().get(
                spreadsheetId=file_id, range="A1:J50"
            ).execute()
            values = sheet.get("values", [])
            return _clean_text("\n".join([", ".join(row) for row in values]))
        
        # Google Slides
        if mime_type == "application/vnd.google-apps.presentation":
            pres = slides_service.presentations().get(presentationId=file_id).execute()
            text = []
            for slide in pres.get("slides", []):
                for el in slide.get("pageElements", []):
                    shape = el.get("shape", {})
                    if "text" in shape:
                        for te in shape["text"].get("textElements", []):
                            if "textRun" in te and "content" in te["textRun"]:
                                text.append(te["textRun"]["content"])
            return _clean_text("".join(text))
        
        # PDF
        if mime_type == "application/pdf":
            if PyPDF2 is None:
                return None
            data = download_file(drive_service, file_dict)
            try:
                reader = PyPDF2.PdfReader(io.BytesIO(data))
                text = []
                for page in reader.pages:
                    try:
                        page_text = page.extract_text()
                        if page_text:
                            text.append(page_text)
                    except Exception:
                        continue
                return _clean_text(" ".join(text))
            except Exception:
                return None
        
        # DOCX
        if mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            if docx is None:
                return None
            data = download_file(drive_service, file_dict)
            try:
                doc = docx.Document(io.BytesIO(data))
                return _clean_text("\n".join([p.text for p in doc.paragraphs]))
            except Exception:
                return None
        
        # PPTX
        if mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            if Presentation is None:
                return None
            data = download_file(drive_service, file_dict)
            try:
                prs = Presentation(io.BytesIO(data))
                text = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            text.append(shape.text)
                return _clean_text("\n".join(text))
            except Exception:
                return None
        
        # Plain text
        if mime_type == "text/plain":
            data = download_file(drive_service, file_dict)
            if _likely_binary(data):
                return None
            try:
                return _clean_text(data.decode("utf-8", errors="ignore"))
            except Exception:
                return None
        
        # Fallback
        raw = download_file(drive_service, file_dict)
        if _likely_binary(raw):
            return None
        try:
            return _clean_text(raw.decode("utf-8", errors="ignore"))
        except Exception:
            return None
            
    except Exception as e:
        print(f"‚ö†Ô∏è Error extracting content: {e}")
        return None

# --------------------
# AI Text Operations
# --------------------
def summarize_text(text):
    if not text or not text.strip():
        return "‚ö†Ô∏è No content to summarize."
    
    snippet = text.strip()
    if len(snippet) > 12000:
        snippet = snippet[:12000]
    
    try:
        resp = co.summarize(
            text=snippet,
            length="short",
            format="paragraph",
            model="summarize-xlarge"
        )
        return (resp.summary or "").strip()
    except Exception as e:
        return f"‚ö†Ô∏è Summarization failed: {e}"

# ===== IMPROVED EMAIL DRAFTING FUNCTIONS =====
def generate_email_draft(context: str):
    """
    Generate a professional email draft based on the provided context.
    Now with better prompting and error handling.
    """
    if not context or len(context.strip()) < 5:
        return None, "Please provide more details about what email you want to draft."
    
    prompt = f"""You are a professional email assistant. Generate a complete, ready-to-send email based on this request:

Request: {context}

Create a JSON response with:
1. "subject": A clear, concise subject line (max 10 words)
2. "body": The complete email body with:
   - Appropriate greeting
   - Clear main message
   - Professional closing
   - Signature placeholder: [Your Name]

The email should be professional, clear, and ready to send immediately.

Return ONLY valid JSON, no other text:"""
    
    try:
        resp = co.chat(model="command-r-plus", message=prompt, temperature=0.3)
        text = (resp.text or "").strip()
        
        # Clean markdown
        text = re.sub(r"^```(?:json)?\s*|\s*```$", "", text, flags=re.IGNORECASE | re.DOTALL).strip()
        
        # Try to parse JSON
        try:
            data = json.loads(text)
            subject = data.get("subject", "").strip()
            body = data.get("body", "").strip()
            
            if not subject:
                # Generate subject from first line of body
                first_line = body.split('\n')[0] if body else ""
                subject = first_line[:50] if first_line else "Email from Workspace Agent"
            
            return subject, body
            
        except json.JSONDecodeError:
            # Fallback: try to extract subject and body from text
            lines = text.split('\n')
            subject = ""
            body_parts = []
            
            for line in lines:
                if line.lower().startswith("subject:") or line.lower().startswith("subject:"):
                    subject = line.split(":", 1)[1].strip()
                elif line.strip() and not line.lower().startswith(("json", "```")):
                    body_parts.append(line.strip())
            
            if not subject and body_parts:
                subject = body_parts[0][:50]
            
            body = "\n".join(body_parts) if body_parts else text
            return subject, body
            
    except Exception as e:
        return None, f"‚ö†Ô∏è Drafting failed: {e}"
    
    return None, "‚ö†Ô∏è Could not generate email draft."

def refine_email_draft(instruction: str, current_subject: str, current_body: str):
    """
    Refine the current draft based on natural language instruction.
    """
    if not current_body:
        return current_subject, "‚ö†Ô∏è No draft exists to refine."
    
    prompt = f"""You are a professional email assistant. Refine this email draft according to the instruction.

Current Email:
Subject: {current_subject or '(No subject)'}
Body: {current_body}

Instruction: {instruction}

Return ONLY a JSON object with "subject" and "body" containing the refined email.
Keep the same general content but adjust the tone/length/style as requested."""
    
    try:
        resp = co.chat(model="command-r-plus", message=prompt, temperature=0.2)
        text = (resp.text or "").strip()
        
        # Clean markdown
        text = re.sub(r"^```(?:json)?\s*|\s*```$", "", text, flags=re.IGNORECASE | re.DOTALL).strip()
        
        try:
            data = json.loads(text)
            subject = data.get("subject", "").strip() or current_subject
            body = data.get("body", "").strip() or current_body
            return subject, body
        except json.JSONDecodeError:
            # If parsing fails, keep original with note
            return current_subject, current_body + f"\n\n[Refinement attempt failed: using original]"
            
    except Exception as e:
        return current_subject, current_body + f"\n\n[Refinement failed: {e}]"

# --------------------
# Email Sending
# --------------------
def _make_attachment_part(filename: str, file_data: bytes):
    part = MIMEBase("application", "octet-stream")
    part.set_payload(file_data)
    encoders.encode_base64(part)
    part.add_header("Content-Disposition", f'attachment; filename="{filename}"')
    return part

def send_email_with_attachments(gmail_service, recipients, subject, body_text, attachments=None):
    """Send email with optional attachments."""
    if not recipients:
        raise ValueError("No recipients provided.")
    
    if attachments:
        message = MIMEMultipart()
        message["to"] = ", ".join(recipients)
        message["subject"] = subject
        message.attach(MIMEText(body_text, "plain"))
        
        for filename, filedata in attachments:
            part = _make_attachment_part(filename, filedata)
            message.attach(part)
    else:
        message = MIMEText(body_text, "plain")
        message["to"] = ", ".join(recipients)
        message["subject"] = subject
    
    raw_message = base64.urlsafe_b64encode(message.as_bytes()).decode()
    gmail_service.users().messages().send(userId="me", body={"raw": raw_message}).execute()

# --------------------
# Interactive Draft Creation
# --------------------
def interactive_draft_prompt():
    """Guide user through creating a detailed email draft."""
    print("\nüìù Let's create your email draft. Please provide:")
    
    purpose = input("Purpose of this email (e.g., sick leave, meeting request, etc.): ").strip()
    if not purpose:
        return None, None
    
    recipient_type = input("Who is this for? (manager, colleague, client, team, etc.): ").strip()
    details = input("Key details/dates/reason: ").strip()
    tone = input("Tone? (formal, casual, urgent, polite) [formal]: ").strip() or "formal"
    
    context = f"Email purpose: {purpose}\nRecipient: {recipient_type}\nDetails: {details}\nTone: {tone}"
    return generate_email_draft(context)

# --------------------
# Main Agent
# --------------------
def print_help():
    print("\n" + "="*80)
    print("ü§ñ WORKSAPCE AGENT - AVAILABLE COMMANDS")
    print("="*80)
    print("\nüìÅ FILE OPERATIONS:")
    print("  ‚Ä¢ list all files")
    print("  ‚Ä¢ search [keyword]")
    print("  ‚Ä¢ summarize [file name]")
    print("  ‚Ä¢ summarize file id [file-id]")
    
    print("\nüìé SEND FILES:")
    print("  ‚Ä¢ send [file name] to [email]")
    print("  ‚Ä¢ send summary of [file] to [email]")
    
    print("\nüìß DRAFT EMAIL (IMPROVED):")
    print("  ‚Ä¢ draft [your complete request]")
    print("  ‚Ä¢ draft - Start with 'draft' and full context")
    print("  ‚Ä¢ Example: draft sick leave request for 2 days starting Monday to my manager")
    print("  ‚Ä¢ Example: compose meeting request about Q4 planning for next Tuesday")
    
    print("\n‚úèÔ∏è MANAGE DRAFTS:")
    print("  ‚Ä¢ make it more formal/polite/professional")
    print("  ‚Ä¢ shorten it / expand it / rewrite")
    print("  ‚Ä¢ show draft")
    print("  ‚Ä¢ clear draft")
    print("  ‚Ä¢ send draft to [email(s)]")
    
    print("\nüÜò OTHER:")
    print("  ‚Ä¢ help")
    print("  ‚Ä¢ exit")
    print("="*80 + "\n")

def _pick_best_match(hits, desired_name: str):
    if not hits:
        return None
    if desired_name:
        for f in hits:
            if f.get("name", "").lower() == desired_name.lower():
                return f
    return hits[0]

def main():
    global current_draft
    
    try:
        drive_service, sheets_service, docs_service, slides_service, gmail_service = get_services()
    except Exception as e:
        print(f"‚ùå Failed to initialize services: {e}")
        return
    
    print("\n" + "="*80)
    print("‚úÖ AGENT IS READY - Type 'help' for commands or 'exit' to quit")
    print("="*80)
    
    while True:
        try:
            command = input("\nüéØ Enter command: ").strip()
            if not command:
                continue
            
            parsed = parse_command(command)
            
            if not parsed or "action" not in parsed:
                print("‚ùå Command not recognized. Try 'help'.")
                continue
            
            action = parsed["action"]
            
            # ---------- EXIT / HELP ----------
            if action == "exit":
                print("üëã Goodbye!")
                break
            if action == "help":
                print_help()
                continue
            
            # ===== IMPROVED DRAFT EMAIL =====
            if action == "draft_email":
                context = parsed.get("text", "").strip()
                
                if not context or len(context.split()) <= 2:
                    print("\n‚ÑπÔ∏è I need more details to draft a proper email.")
                    subject, body = interactive_draft_prompt()
                else:
                    print(f"\nüìù Drafting email based on: '{context}'")
                    subject, body = generate_email_draft(context)
                
                if not body or body.startswith("‚ö†Ô∏è"):
                    print(f"\n{body if body else '‚ùå Failed to create draft.'}")
                    continue
                
                # Store draft
                current_draft = {
                    "subject": subject or "Drafted Email",
                    "body": body,
                    "recipients": parse_recipients(context + " " + command),
                    "context": context
                }
                
                print("\n" + "="*80)
                print("üìß DRAFT CREATED")
                print("="*80)
                print(f"Subject: {current_draft['subject']}\n")
                print(current_draft['body'])
                print("\n" + "="*80)
                print("Commands: 'make it more formal', 'shorten it', 'send draft to email@domain.com'")
                print("="*80)
                continue
            
            # ===== REFINE DRAFT =====
            if action == "refine_draft":
                if not current_draft["body"]:
                    print("‚ö†Ô∏è No draft exists. Create one first with 'draft ...'")
                    continue
                
                instruction = parsed.get("instruction", command)
                print(f"\n‚úèÔ∏è Refining draft: '{instruction}'")
                
                subject, body = refine_email_draft(
                    instruction, 
                    current_draft["subject"], 
                    current_draft["body"]
                )
                
                current_draft["subject"] = subject or current_draft["subject"]
                current_draft["body"] = body or current_draft["body"]
                
                print("\n" + "="*80)
                print("üìß UPDATED DRAFT")
                print("="*80)
                print(f"Subject: {current_draft['subject']}\n")
                print(current_draft['body'])
                print("="*80)
                continue
            
            # ===== SHOW DRAFT =====
            if action == "show_draft":
                if not current_draft["body"]:
                    print("‚ö†Ô∏è No draft currently stored.")
                else:
                    print("\n" + "="*80)
                    print("üìß CURRENT DRAFT")
                    print("="*80)
                    print(f"Subject: {current_draft['subject']}")
                    print(f"Recipients: {', '.join(current_draft['recipients']) if current_draft['recipients'] else 'Not specified'}")
                    print("\n" + current_draft['body'])
                    print("="*80)
                continue
            
            # ===== CLEAR DRAFT =====
            if action == "clear_draft":
                current_draft = {
                    "subject": None,
                    "body": None,
                    "recipients": [],
                    "context": ""
                }
                print("‚úÖ Draft cleared.")
                continue
            
            # ===== SEND DRAFT =====
            if action == "send_draft":
                if not current_draft["body"]:
                    print("‚ö†Ô∏è No draft exists to send.")
                    continue
                
                recipients = parsed.get("email") or parse_recipients(command) or current_draft["recipients"]
                
                if not recipients:
                    print("üìß Who should I send this draft to? (comma-separated emails)")
                    rline = input("Recipients: ").strip()
                    recipients = parse_recipients(rline)
                
                if not recipients:
                    print("‚ùå No valid recipients provided.")
                    continue
                
                try:
                    send_email_with_attachments(
                        gmail_service,
                        recipients,
                        current_draft["subject"] or "Drafted Email",
                        current_draft["body"]
                    )
                    print(f"‚úÖ Email sent successfully to: {', '.join(recipients)}")
                    
                    # Clear draft after successful send
                    current_draft = {
                        "subject": None,
                        "body": None,
                        "recipients": [],
                        "context": ""
                    }
                    print("‚úÖ Draft cleared after sending.")
                    
                except Exception as e:
                    print(f"‚ùå Failed to send email: {e}")
                continue
            
            # ---------- LIST FILES ----------
            if action == "list_files":
                list_all_files(drive_service)
                continue
            
            # ---------- SEARCH FILES ----------
            if action == "search_files":
                keyword = parsed.get("keyword")
                keywords = parsed.get("keywords")
                all_hits = []
                
                if keyword:
                    all_hits.extend(search_files(drive_service, keyword))
                if isinstance(keywords, list):
                    for kw in keywords:
                        all_hits.extend(search_files(drive_service, kw))
                
                all_hits = _unique_by_id(all_hits)
                
                if not all_hits:
                    term = keyword if keyword else ", ".join(keywords or [])
                    print(f"‚ùå No files found for '{term}'.")
                else:
                    _print_files(all_hits)
                continue
            
            # ---------- SEND FILES ----------
            if action == "send_file":
                recipients = parse_recipients(command)
                if not recipients and parsed.get("email"):
                    recipients = [parsed.get("email")]
                
                if not recipients:
                    print("‚ùå Missing recipient email(s).")
                    continue
                
                names = parsed.get("file_name", [])
                if isinstance(names, str):
                    names = [names]
                
                if not names:
                    print("‚ùå No file names specified.")
                    continue
                
                matched_all = []
                for name in names:
                    matched = search_files(drive_service, name)
                    matched_all.extend(matched)
                
                matched_all = _unique_by_id(matched_all)
                
                if not matched_all:
                    print("‚ö†Ô∏è No files matched your request.")
                    continue
                
                wants_summary = bool(re.search(r'\bsummary\b|\bsummar(y|ise|ize)\b', command.lower()))
                
                if wants_summary:
                    # Send summaries
                    all_summaries = []
                    for f in matched_all[:3]:  # Limit to 3 files for summaries
                        actual = _deref_shortcut(drive_service, f)
                        file_name = actual.get("name", "(unknown)")
                        content = extract_file_content(drive_service, sheets_service, docs_service, slides_service, actual)
                        
                        if not content:
                            all_summaries.append(f"‚ö†Ô∏è Could not extract content for '{file_name}'")
                            continue
                        
                        summary = summarize_text(content)
                        all_summaries.append(f"--- {file_name} ---\n{summary}\n")
                    
                    combined = "\n\n".join(all_summaries).strip()
                    subject = "File Summaries from Workspace Agent"
                    body = f"Here are the summaries you requested:\n\n{combined}"
                    
                    try:
                        send_email_with_attachments(gmail_service, recipients, subject, body)
                        print(f"‚úÖ Sent summaries to: {', '.join(recipients)}")
                    except Exception as e:
                        print(f"‚ùå Failed to send summaries: {e}")
                        
                else:
                    # Send files as attachments
                    attachments = []
                    failed = []
                    
                    for f in matched_all:
                        try:
                            actual = _deref_shortcut(drive_service, f)
                            data = download_file(drive_service, actual)
                            attachments.append((actual.get("name", "file"), data))
                        except Exception as e:
                            failed.append((f.get("name", "(unknown)"), str(e)))
                    
                    if attachments:
                        subject = "Files from Workspace Agent"
                        body = f"Attached {len(attachments)} file(s) as requested."
                        
                        try:
                            send_email_with_attachments(gmail_service, recipients, subject, body, attachments)
                            print(f"‚úÖ Sent {len(attachments)} file(s) to: {', '.join(recipients)}")
                            
                            if failed:
                                print("‚ö†Ô∏è Some files failed to attach:")
                                for name, err in failed:
                                    print(f"   - {name}: {err}")
                                    
                        except Exception as e:
                            print(f"‚ùå Failed to send attachments: {e}")
                    else:
                        print("‚ö†Ô∏è No files could be attached.")
                        for name, err in failed:
                            print(f"   - {name}: {err}")
                
                continue
            
            # ---------- SUMMARIZE FILE ----------
            if action == "summarize_file":
                file_dict = None
                file_id = parsed.get("file_id")
                file_name = parsed.get("file_name")
                
                try:
                    if file_id:
                        file_dict = get_file_by_id(drive_service, file_id)
                    elif file_name:
                        hits = search_files(drive_service, file_name)
                        if hits:
                            file_dict = _pick_best_match(hits, file_name)
                    
                    if not file_dict:
                        print("‚ùå Could not find the file to summarize.")
                        continue
                    
                    file_dict = _deref_shortcut(drive_service, file_dict)
                    content = extract_file_content(drive_service, sheets_service, docs_service, slides_service, file_dict)
                    
                    if not content:
                        mime = file_dict.get("mimeType", "")
                        if "presentation" in mime and Presentation is None:
                            print("‚ö†Ô∏è PPTX summarization requires 'python-pptx'. Install with: pip install python-pptx")
                        else:
                            print("‚ö†Ô∏è Could not extract text content from this file.")
                        continue
                    
                    summary = summarize_text(content)
                    print(f"\nüìå SUMMARY: {file_dict.get('name', '(unknown)')}")
                    print("-" * 80)
                    print(summary)
                    print("-" * 80)
                    
                except Exception as e:
                    print(f"‚ùå Summarization failed: {e}")
                continue
            
            print("‚ùå Command not fully implemented. Try 'help'.")
            
        except KeyboardInterrupt:
            print("\n\nüëã Goodbye!")
            break
        except Exception as e:
            print(f"‚ùå Error: {e}")
            print("Try again or type 'help'.")

if __name__ == "__main__":
    main()